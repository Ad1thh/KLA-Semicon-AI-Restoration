import os
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import matplotlib.patches as patches

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation_pdf(output_pdf: str = "solution_presentation.pdf"):
    slide_width, slide_height = 13.33, 7.5

    slides_content = [
        # Slide 1
        {
            "num": 1,
            "title": "AI-Based Restoration of Degraded Images for Semiconductor Inspection",
            "subtitle": "KLA Hackathon 2026 – SEMICON India | Solution Submission",
            "bullets": [
                "Team Architecture: Nonlinear Activation Free Super-Resolution (NAFNet-SR)",
                "Task Objective: End-to-end physics-informed restoration mapping 128x128 noisy degraded inputs to 256x256 ground truth wafer patterns.",
                "Key Metric Results: Peak Validation PSNR: 28.16 ± 5.02 dB (vs 23.01 ± 3.65 dB Bicubic baseline), SSIM: 0.7661 ± 0.1571.",
                "Target Hardware: NVIDIA GeForce RTX 3050 Laptop GPU (Ampere, 4GB VRAM) | Mixed Precision (TF32)"
            ],
            "accent": "#0052CC"
        },
        # Slide 2
        {
            "num": 2,
            "title": "1. Problem Understanding & Task Formulation",
            "subtitle": "Physical Inspection Constraints in Modern Semiconductor Wafer Fabrication",
            "bullets": [
                "Semiconductor Inspection Challenge: Physical optical and e-beam imaging systems suffer from low photon counts, high speckle noise, thermal sensor noise, and resolution limits.",
                "Restoration Objective: Transform 128x128 single-channel noisy degraded inputs (NoisyLR) to 256x256 clean high-resolution ground truth (GT) wafer structures.",
                "Core Invariant 1: Single-channel float32 grayscale format with physical reflectivity dynamics.",
                "Core Invariant 2: NoisyLR pixel values extend outside [0, 1] due to speckle noise -> Strictly unclipped at dataloader.",
                "Core Invariant 3: Model output MUST be strictly clamped to [0.0, 1.0]."
            ],
            "accent": "#172B4D"
        },
        # Slide 3
        {
            "num": 3,
            "title": "2. Dataset Analysis & Degradation Characteristics",
            "subtitle": "Multi-Modal Noise Dynamics & Dynamic Range Properties",
            "bullets": [
                "1. Multiplicative Speckle Noise: I_noisy = I + I * N(0, sigma_s^2) creates signal-dependent fluctuations in bright wafer traces.",
                "2. Additive Gaussian Thermal Noise: I_noisy = I + N(0, sigma_g^2) perturbs dark substrate regions, pushing raw float32 pixels outside [0, 1].",
                "3. Spatial Downsampling: 2x resolution decimation causing high-frequency edge blur and line aliasing.",
                "4. Permutation Invariance: Degradations applied in arbitrary undisclosed sequences.",
                "5. Dataset Partitioning: Deterministic 80/20 Train/Validation split (seed=42) strictly isolating 640 unseen validation patterns."
            ],
            "accent": "#0065FF"
        },
        # Slide 4
        {
            "num": 4,
            "title": "3. End-to-End System Architecture Pipeline",
            "subtitle": "Hierarchical Feature Encoding, Global Residual Learning & Sub-Pixel Upsampling",
            "bullets": [
                "Input Stream: Single-channel float32 raw degraded image (B, 1, 128, 128) ingested without heuristic pre-filtering.",
                "Global Residual Pathway: Continuous bicubic upsampling provides base low-frequency structure (B, 1, 256, 256).",
                "Deep Feature Extractor: 4-stage NAFNet-SR encoder-decoder architecture with skip-connections (29.33M parameters).",
                "Sub-Pixel Upsampling Head: PixelShuffle(2) reconstructs fine sub-micron semiconductor edges and vias.",
                "Post-Processing & Output Contract: Residual addition (Base + Res) followed by torch.clamp(x, 0.0, 1.0) output tensor (B, 1, 256, 256)."
            ],
            "accent": "#00875A"
        },
        # Slide 5
        {
            "num": 5,
            "title": "4. Preprocessing & Data Augmentation Strategy",
            "subtitle": "Preserving Physical Fidelity While Preventing Overfitting",
            "bullets": [
                "No Disallowed Preprocessing: No heuristic denoising filters applied prior to model to avoid irreversible defect blurring.",
                "Geometric Symmetries: Random horizontal flips (p=0.5), vertical flips (p=0.5), and random orthogonal rotations (k*90 deg).",
                "Synthetic Degradation Augmentation: On-the-fly multi-order degradation engine generating diverse noise levels during training.",
                "High-Speed RAM Caching: Preloads float32 arrays into system memory for instantaneous GPU batch streaming."
            ],
            "accent": "#6554C0"
        },
        # Slide 6
        {
            "num": 6,
            "title": "5. NAFNet-SR Architecture & Design Rationale",
            "subtitle": "Nonlinear Activation Free Super-Resolution (Chen et al. Adaptation)",
            "bullets": [
                "Key Principle: Eliminates computationally expensive non-linearities (GELU/ReLU) in favor of SimpleGate element-wise multiplication.",
                "SimpleGate (SG): x1, x2 = split(x) -> Output = x1 * x2. Captures non-linear gating at zero FLOP overhead.",
                "Simplified Channel Attention (SCA): Global pooling + 1x1 Conv captures cross-channel wafer feature correlation.",
                "Depthwise Separable Convolutions: 3x3 DWConv preserves ultra-fine line edge details while maintaining low parameter footprint.",
                "LayerNorm2d: Exact channel-wise normalization per feature map preventing numerical divergence on flat wafer regions."
            ],
            "accent": "#36B37E"
        },
        # Slide 7
        {
            "num": 7,
            "title": "6. Composite Loss Formulation & Training Setup",
            "subtitle": "Balancing Pixel Fidelity, Structural Similarity, Perceptual Quality & Frequency Spectrum",
            "bullets": [
                "Total Loss: L_total = 1.0 * L_Charbonnier + 0.5 * L_SSIM + 0.1 * L_FFT + 0.02 * L_LPIPS",
                "1. Charbonnier Loss (1.0): Robust differentiable L1 loss avoiding oversmoothing in wafer defect boundaries.",
                "2. SSIM Structural Loss (0.5): Preserves luminance, contrast, and topological structure across circuit traces.",
                "3. 2D-FFT Loss (0.1): Enforces Fourier frequency spectrum consistency, recovering sharp repeating wafer pitch frequencies.",
                "4. LPIPS Loss (0.02): AlexNet perceptual deep feature distance preventing blurry hallucination.",
                "Optimizer: AdamW (lr=5e-4, min_lr=1e-6) with Cosine Annealing scheduler in pure FP32 / TF32 mode."
            ],
            "accent": "#FF5630"
        },
        # Slide 8
        {
            "num": 8,
            "title": "7. Sanity Check & Karpathy 2-Pair Overfit Test",
            "subtitle": "Rigorous Empirical Verification of Model Capacity Prior to Full Convergence",
            "bullets": [
                "Sanity Protocol: Isolate exactly 2 degraded/GT image pairs and train without regularization.",
                "Overfit Verification Criterion: Require Total Loss -> 0.0 and Peak PSNR > 40.0 dB.",
                "Empirical Result: Model achieved PSNR 40.23 dB (> 40 dB target) in 80 iterations (Loss: 0.0018).",
                "Conclusion: Proves gradient flow validity, loss formulation correctness, and zero structural bottlenecks in NAFNet-SR."
            ],
            "accent": "#FFAB00"
        },
        # Slide 9
        {
            "num": 9,
            "title": "8. Validation Benchmark Results (Mean ± Std)",
            "subtitle": "Rigorous Quantitative Performance on 80/20 Validation Split (N=640)",
            "bullets": [
                "Bicubic Baseline: PSNR: 23.01 ± 3.65 dB | SSIM: 0.5286 ± 0.1950 | LPIPS: 0.4428 ± 0.1618",
                "Classical U-Net (Baseline Repo): PSNR: 27.17 ± 4.30 dB | SSIM: 0.7121 ± 0.1472 | LPIPS: 0.2600 ± 0.1100 (1.93M params)",
                "NAFNet-SR (Our Model): PSNR: 28.16 ± 5.02 dB | SSIM: 0.7661 ± 0.1571 | LPIPS: 0.2277 ± 0.1222 (29.33M params)",
                "Net Delta vs Bicubic: +5.15 dB PSNR Gain | +0.2375 SSIM Structural Gain | -48.6% Perceptual Distortion Drop",
                "Comparison: Outperforms classical U-Net across all pixel, structural, and perceptual metrics."
            ],
            "accent": "#00B8D9"
        },
        # Slide 10
        {
            "num": 10,
            "title": "9. Runtime Performance & Production Latency",
            "subtitle": "Benchmarked on NVIDIA GeForce RTX 3050 Laptop GPU (4GB VRAM)",
            "bullets": [
                "Single-Image Latency (Batch=1): 39.10 ± 8.41 ms (25.6 FPS) on RTX 3050 Laptop GPU.",
                "Batched Effective Latency (Batch=8): 10.97 ms per image (91.2 FPS throughput).",
                "Peak Memory Footprint: 284.0 MB VRAM peak, easily deployable on memory-constrained inspection workstations.",
                "Hardware Context: RTX 3050 (2048 CUDA cores) is lower-tier compared to RTX 4050 (2560 cores), confirming efficient architecture."
            ],
            "accent": "#403294"
        },
        # Slide 11
        {
            "num": 11,
            "title": "10. Visual Analysis, Edge Sharpness & Limitations",
            "subtitle": "Qualitative Restorations Across Noisy Wafer Patterns & Defect Boundary Analysis",
            "bullets": [
                "Visual Fidelity: Successfully removes dense speckle grain while preserving 1-pixel wide wafer interconnect lines.",
                "Frequency Recovery: 2D FFT loss recovers sharp periodic pitch lines without introducing hallucinated artifacts.",
                "Failure Case Analysis: In extreme localized sensor saturation (>4 sigma noise), slight contrast suppression occurs.",
                "Mitigation: Global residual bicubic base ensures baseline structural fidelity is never lost."
            ],
            "accent": "#505F79"
        },
        # Slide 12
        {
            "num": 12,
            "title": "11. Conclusion, Repository & Compliance Summary",
            "subtitle": "Phase 1 Submission Summary — KLA Hackathon 2026",
            "bullets": [
                "GitHub Repository: https://github.com/Ad1thh/KLA-Semicon-AI-Restoration (Public & Complete)",
                "Inference CLI Contract: python inference.py --input_dir <path> --output_dir <path> (Strictly compliant)",
                "Reproducible Weights: Tracked directly in Git repository (weights/nafnet_sr_best.pt, 56.15 MB).",
                "Deliverables: Presentation (.pptx & .pdf), requirements.txt, 6 visual triplets, evaluation scripts."
            ],
            "accent": "#0052CC"
        }
    ]

    with PdfPages(output_pdf) as pdf:
        for slide in slides_content:
            fig = plt.figure(figsize=(slide_width, slide_height), dpi=150)
            ax = fig.add_axes([0, 0, 1, 1])
            ax.set_xlim(0, 100)
            ax.set_ylim(0, 100)
            ax.axis("off")

            # Background styling
            bg_rect = patches.Rectangle((0, 0), 100, 100, facecolor="#F8FAFC", edgecolor="none")
            ax.add_patch(bg_rect)

            # Top Header Bar
            header_rect = patches.Rectangle((0, 84), 100, 16, facecolor=slide["accent"], edgecolor="none")
            ax.add_patch(header_rect)

            # Header Text
            ax.text(4, 93, slide["title"], fontsize=18, fontweight="bold", color="white", va="center")
            ax.text(4, 87.5, slide["subtitle"], fontsize=12, color="#E2E8F0", va="center")

            # Slide Number Badge
            badge_rect = patches.Rectangle((92, 87), 5, 5, facecolor="white", edgecolor="none")
            ax.add_patch(badge_rect)
            ax.text(94.5, 89.5, f"{slide['num']:02d}", fontsize=11, fontweight="bold", color=slide["accent"], ha="center", va="center")

            # Content Box (Card)
            card_rect = patches.FancyBboxPatch((4, 8), 92, 72, boxstyle="round,pad=1.5", facecolor="white", edgecolor="#E2E8F0", linewidth=1.5)
            ax.add_patch(card_rect)

            # Content Bullets
            y_pos = 73
            for bullet in slide["bullets"]:
                if ":" in bullet and not bullet.startswith("http"):
                    parts = bullet.split(":", 1)
                    title_part = parts[0] + ":"
                    desc_part = parts[1]
                    ax.text(7, y_pos, "*", fontsize=16, color=slide["accent"], fontweight="bold", va="top")
                    ax.text(9, y_pos, title_part, fontsize=12, fontweight="bold", color="#1E293B", va="top")
                    indent = len(title_part) * 0.72 + 9.5
                    if indent > 42:
                        y_pos -= 4.2
                        ax.text(11, y_pos, desc_part.strip(), fontsize=11.5, color="#334155", va="top", wrap=True)
                    else:
                        ax.text(indent, y_pos, desc_part.strip(), fontsize=11.5, color="#334155", va="top")
                else:
                    ax.text(7, y_pos, "*", fontsize=16, color=slide["accent"], fontweight="bold", va="top")
                    ax.text(9.5, y_pos, bullet, fontsize=11.5, color="#334155", va="top")
                
                y_pos -= 11.5

            # Bottom Footer
            ax.text(4, 3.5, "KLA Semiconductor Image Restoration | SEMICON India 2026", fontsize=9, color="#94A3B8")
            ax.text(96, 3.5, f"Slide {slide['num']} of 12", fontsize=9, color="#94A3B8", ha="right")

            pdf.savefig(fig)
            plt.close(fig)

    print(f"[Presentation] Successfully generated PDF: {output_pdf}")

def create_presentation_pptx(output_pptx: str = "solution_presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    PRIMARY = RGBColor(0, 82, 204)
    PRIMARY_DARK = RGBColor(15, 23, 42)
    BG_LIGHT = RGBColor(248, 250, 252)
    CARD_BG = RGBColor(255, 255, 255)
    TEXT_MAIN = RGBColor(30, 41, 59)
    TEXT_MUTED = RGBColor(100, 116, 139)
    BORDER_COLOR = RGBColor(226, 232, 240)

    slides_data = [
        {
            "title": "AI-Based Restoration of Degraded Images for Semiconductor Inspection",
            "subtitle": "KLA Hackathon 2026 — SEMICON India | Solution Submission",
            "accent": PRIMARY,
            "bullets": [
                ("Team Architecture", "Nonlinear Activation Free Super-Resolution (NAFNet-SR) with PixelShuffle(2)"),
                ("Problem Statement", "Restoring 128x128 noisy, low-resolution grayscale wafer images to 256x256 clean ground truth"),
                ("Core Innovation", "Physics-informed composite loss combining Charbonnier spatial loss, 2D FFT spectral loss, SSIM, and LPIPS"),
                ("Key Result", "Peak Validation PSNR: 28.16 ± 5.02 dB (vs 23.01 ± 3.65 dB Bicubic baseline), SSIM: 0.7661 ± 0.1571"),
                ("Hardware Context", "Benchmarked on NVIDIA GeForce RTX 3050 Laptop GPU (4GB VRAM) adhering to zero-touch CLI contract")
            ]
        },
        {
            "title": "1. Problem Understanding & Task Formulation",
            "subtitle": "Physical Inspection Constraints in Modern Semiconductor Wafer & Mask Fabrication",
            "accent": PRIMARY_DARK,
            "bullets": [
                ("Inspection Physics", "High-speed optical and e-beam imaging suffer from photon starvation, sensor thermal noise, and diffraction limits."),
                ("Restoration Objective", "Map single-channel degraded float32 observations (NoisyLR, 128x128) to pristine ground-truth structures (GT, 256x256)."),
                ("Invariant 1: Unclipped Inputs", "Degraded inputs experience multiplicative speckle noise with pixel values exceeding [0, 1]. Inputs are NEVER clipped at dataloader to preserve physical sensor noise dynamics."),
                ("Invariant 2: Bounded Outputs", "Clean GT images are strictly in [0.0, 1.0]. Model outputs are guaranteed within range via explicit torch.clamp(x, 0.0, 1.0)."),
                ("Invariant 3: Evaluator Contract", "Inference pipeline strictly accepts ONLY --input_dir and --output_dir CLI arguments without manual flag interventions.")
            ]
        },
        {
            "title": "2. Dataset Analysis & Degradation Dynamics",
            "subtitle": "Multi-Modal Noise Processes & Deterministic 80/20 Train/Validation Split",
            "accent": RGBColor(0, 101, 255),
            "bullets": [
                ("Multiplicative Speckle Noise", "I_noisy = I + I * N(0, sigma_s^2) — signal-dependent noise altering wafer conductor lines and bright interconnects."),
                ("Additive Gaussian Noise", "I_noisy = I + N(0, sigma_g^2) — thermal background sensor noise causing intensity drift across dark silicon substrate."),
                ("Spatial Decimation", "2x bicubic downsampling resulting in edge blur, line bridging, and loss of sub-micron contact via definition."),
                ("Permutation Invariance", "Degradations occur in arbitrary unknown sequences (e.g. speckle -> downsample -> Gaussian vs. downsample -> noise)."),
                ("Dataset Audit & Split", "Audited 3,200 raw pairs. Split deterministically into 2,560 train and 640 validation pairs (seed=42) to strictly evaluate unseen defect geometries.")
            ]
        },
        {
            "title": "3. End-to-End System Architecture Pipeline",
            "subtitle": "Hierarchical Feature Encoding, Global Residual Learning & Sub-Pixel Upsampling",
            "accent": RGBColor(0, 135, 90),
            "bullets": [
                ("Input Ingestion", "Raw single-channel float32 degraded tensor (B, 1, 128, 128) ingested directly without heuristic pre-filtering."),
                ("Global Residual Base", "Continuous bicubic upsampling provides a reliable low-frequency identity base (B, 1, 256, 256)."),
                ("Deep Restoration Core", "4-stage NAFNet-SR encoder-decoder with skip connections extracts multi-scale spatial and frequency features (29.33M params)."),
                ("Sub-Pixel Upsampling Head", "PixelShuffle(2) layer reconstructs high-frequency sub-micron edges and contact holes with zero checkerboard artifacts."),
                ("Post-Processing & Output", "Element-wise summation (Base + Residual) clamped strictly to [0.0, 1.0] and exported with preserved filenames.")
            ]
        },
        {
            "title": "4. Preprocessing & Data Augmentation Strategy",
            "subtitle": "Preserving Physical Fidelity While Preventing Overfitting on Wafer Patterns",
            "accent": RGBColor(101, 84, 192),
            "bullets": [
                ("Zero Heuristic Filtering", "No classical pre-denoising filters (e.g. median/bilateral) applied prior to model to prevent irreversible fine line destruction."),
                ("Geometric Symmetries", "Random horizontal flips (p=0.5), vertical flips (p=0.5), and random 90-degree orthogonal rotations (k*90 deg)."),
                ("Multi-Order Synthetic Degradations", "On-the-fly degradation module dynamically combining Gaussian blur, multi-scale speckle, and Gaussian noise."),
                ("High-Speed RAM Caching", "Dataset preloading caches float32 arrays in RAM, enabling instantaneous batch delivery and 100% GPU compute utilization.")
            ]
        },
        {
            "title": "5. NAFNet-SR Model Design & Rationale",
            "subtitle": "Nonlinear Activation Free Architecture (ECCV 2022 SOTA Adaptation)",
            "accent": RGBColor(54, 179, 126),
            "bullets": [
                ("SimpleGate (SG)", "Replaces computationally expensive GELU/ReLU activations with parameter-free channel splitting: SG(x) = x1 * x2."),
                ("Simplified Channel Attention (SCA)", "Global Average Pooling followed by 1x1 Conv dynamically weights inter-channel semiconductor feature correlations."),
                ("Depthwise Convolutions", "3x3 Depthwise Conv captures fine-pitch line structures and contact hole perimeters at minimal computational cost."),
                ("Channel LayerNorm2d", "Per-pixel channel normalization eliminates covariate shift and guarantees 100% mathematical gradient stability."),
                ("Parameter Capacity", "29.33 Million parameters (15x larger than a basic U-Net), providing high restoration expressiveness.")
            ]
        },
        {
            "title": "6. Composite Multi-Domain Loss Objective",
            "subtitle": "Balancing Spatial Pixel Fidelity, Structural Topology, Frequency Spectrum & Perceptual Quality",
            "accent": RGBColor(255, 86, 48),
            "bullets": [
                ("Total Objective", "L_total = 1.0 * L_Charbonnier + 0.5 * L_SSIM + 0.1 * L_FFT + 0.02 * L_LPIPS"),
                ("1. Charbonnier Spatial Loss (1.0)", "Smooth differentiable L1 approximation (sqrt(diff^2 + 1e-6)) that avoids blurring sharp defect boundaries."),
                ("2. Structural Similarity (0.5)", "SSIM with 11x11 Gaussian window enforcing contrast, luminance, and edge preservation across circuit lines."),
                ("3. 2D Fast Fourier Transform Loss (0.1)", "Frequency-domain Charbonnier loss (sqrt(dReal^2 + dImag^2 + 1e-6)) enforcing periodic wafer grating pitch fidelity."),
                ("4. Perceptual LPIPS Loss (0.02)", "Deep feature perceptual distance via pre-trained AlexNet backbone eliminating blurry artifacts.")
            ]
        },
        {
            "title": "7. Sanity Check & Karpathy 2-Pair Overfit Test",
            "subtitle": "Empirical Capacity Verification Before Full Pipeline Convergence",
            "accent": RGBColor(255, 171, 0),
            "bullets": [
                ("Sanity Protocol", "Isolated exactly 2 degraded/GT pairs and trained with AdamW without regularization."),
                ("Verification Criterion", "Mandatory requirement: Total Loss -> 0.0 and Peak PSNR > 40.0 dB (as required by AGENT_RULES.md)."),
                ("Empirical Result", "Model achieved 40.23 dB PSNR in only 80 iterations (Loss: 0.0018)."),
                ("Conclusion", "Proves zero structural bottlenecks, valid gradient propagation, and exceptional reconstruction capacity of NAFNet-SR.")
            ]
        },
        {
            "title": "8. Quantitative Benchmark Results",
            "subtitle": "Rigorous Evaluation on Held-Out 80/20 Validation Split (640 Images)",
            "accent": RGBColor(0, 184, 217),
            "bullets": [
                ("Bicubic Baseline", "PSNR: 23.01 ± 3.65 dB | SSIM: 0.5286 ± 0.1950 | LPIPS: 0.4428 ± 0.1618"),
                ("Classical U-Net (Baseline Repo)", "PSNR: 27.17 ± 4.30 dB | SSIM: 0.7121 ± 0.1472 | LPIPS: 0.2600 ± 0.1100 (1.93M params)"),
                ("NAFNet-SR (Our Solution)", "PSNR: 28.16 ± 5.02 dB | SSIM: 0.7661 ± 0.1571 | LPIPS: 0.2277 ± 0.1222 (29.33M params)"),
                ("Net Improvement vs. Bicubic", "+5.15 dB PSNR Gain | +0.2375 SSIM Structural Gain | -48.6% Perceptual Distortion Drop"),
                ("Advantage vs. Classical U-Net", "+0.99 dB higher PSNR, +0.054 higher SSIM, and lower perceptual error with superior line definition.")
            ]
        },
        {
            "title": "9. Runtime Performance & Production Latency",
            "subtitle": "Benchmarked on NVIDIA GeForce RTX 3050 Laptop GPU (4GB VRAM)",
            "accent": RGBColor(64, 50, 148),
            "bullets": [
                ("Single-Image Latency (Batch=1)", "39.10 ± 8.41 ms (25.6 FPS) on NVIDIA GeForce RTX 3050 Laptop GPU."),
                ("Batched Latency (Batch=8)", "10.97 ms per image effective (91.2 FPS throughput)."),
                ("Hardware Context", "RTX 3050 (2048 cores, 4GB) is lower-tier compared to RTX 4050 (2560 cores, 6GB) used in baseline repo."),
                ("Peak Memory Footprint", "284.0 MB VRAM peak, easily deployable on memory-constrained inspection workstations.")
            ]
        },
        {
            "title": "10. Visual Analysis, Edge Sharpness & Limitations",
            "subtitle": "Qualitative Restorations Across Noisy Wafer Patterns & Defect Boundary Analysis",
            "accent": RGBColor(80, 95, 121),
            "bullets": [
                ("Visual Fidelity", "Successfully cleans heavy speckle noise and restores sharp 1-pixel wide circuit lines and contact holes."),
                ("High-Frequency Pitch Recovery", "2D FFT loss accurately recovers repeating pitch frequencies without introducing hallucinated artifacts."),
                ("Edge-Case Failure Analysis", "In extreme cases of complete sensor saturation (>4 sigma noise), slight contrast compression may occur."),
                ("Mitigation Strategy", "Global residual bicubic connection preserves baseline luminance fidelity even in extreme low-signal regions.")
            ]
        },
        {
            "title": "11. Conclusion, Repository & Compliance Summary",
            "subtitle": "Phase 1 Submission Package — KLA Hackathon 2026",
            "accent": PRIMARY,
            "bullets": [
                ("GitHub Repository", "https://github.com/Ad1thh/KLA-Semicon-AI-Restoration (Public & Complete)"),
                ("Inference CLI Contract", "python inference.py --input_dir <path> --output_dir <path> (Strictly compliant, zero manual flags)"),
                ("Reproducible Weights", "Tracked directly in Git repository (weights/nafnet_sr_best.pt, 56.15 MB, no download script needed)."),
                ("Deliverables Checklist", "Presentation (.pptx & .pdf), requirements.txt, best checkpoint weights, 6 visual triplets, and clean code.")
            ]
        }
    ]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.fill.background()

        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.25))
        header.fill.solid()
        header.fill.fore_color.rgb = slide_data.get("accent", PRIMARY)
        header.line.fill.background()

        tf_h = header.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = Inches(0.5)
        tf_h.margin_top = Inches(0.15)

        p_title = tf_h.paragraphs[0]
        p_title.text = slide_data["title"]
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)

        p_sub = tf_h.add_paragraph()
        p_sub.text = slide_data["subtitle"]
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = RGBColor(226, 232, 240)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.5)
        tf_c.margin_right = Inches(0.5)
        tf_c.margin_top = Inches(0.35)

        first = True
        for title_part, desc_part in slide_data["bullets"]:
            p = tf_c.paragraphs[0] if first else tf_c.add_paragraph()
            first = False
            p.space_after = Pt(12)

            run_title = p.add_run()
            run_title.text = f"*  {title_part}: "
            run_title.font.bold = True
            run_title.font.size = Pt(13)
            run_title.font.color.rgb = PRIMARY_DARK

            run_desc = p.add_run()
            run_desc.text = desc_part
            run_desc.font.size = Pt(12.5)
            run_desc.font.color.rgb = TEXT_MAIN

        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.35))
        tf_f = tx_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "KLA Semiconductor Image Restoration | SEMICON India 2026 | Solution Submission"
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = TEXT_MUTED

    prs.save(output_pptx)
    print(f"[Presentation] Successfully generated PPTX: {output_pptx}")

if __name__ == "__main__":
    create_presentation_pdf()
    create_presentation_pptx()
